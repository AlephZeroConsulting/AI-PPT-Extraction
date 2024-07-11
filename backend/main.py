from fastapi import FastAPI, BackgroundTasks
from pydantic import BaseModel
import os
import subprocess
from pathlib import Path
import shutil
import hashlib
from fastapi.responses import JSONResponse, FileResponse
from fastapi.middleware.cors import CORSMiddleware
from pytube import YouTube
import cv2
import base64


import torch
import torch.nn as nn
import torchvision.models as models
import torchvision.transforms as transforms
from PIL import Image
import cv2
import numpy as np
import time
from tqdm import tqdm


from pptx import Presentation
from pptx.util import Inches
import cv2
from glob import glob
import re



os.makedirs('./videos', exist_ok=True)
os.makedirs('./ppts', exist_ok=True)
os.makedirs('./frames', exist_ok=True)

# Check if CUDA (GPU acceleration) is available
device = torch.device('cuda' if torch.cuda.is_available() else 'cpu')

print(f"Using device: {device}")

# Load the pre-trained VGG16 model
model = models.vgg16(pretrained=True)
# Use the output of the second last layer (before the final classification layer) for feature extraction
model.classifier = nn.Sequential(*list(model.classifier.children())[:-1])
model = model.to(device)
model.eval()  # Set the model to evaluation mode

# Define the image transformation
transform = transforms.Compose([
    transforms.Resize((224, 224)),
    transforms.ToTensor(),
    transforms.Normalize(mean=[0.485, 0.456, 0.406], std=[0.229, 0.224, 0.225]),
])

def extract_features(img):
    # Convert the image to PIL format
    img_pil = Image.fromarray(cv2.cvtColor(img, cv2.COLOR_BGR2RGB))
    # Apply the transformations
    img_tensor = transform(img_pil).unsqueeze(0)
    # Extract features
    with torch.no_grad():
        features = model(img_tensor.to(device))
    return features.flatten().cpu().numpy()



def get_image_size(image_path):
    img = cv2.imread(image_path)
    height, width, _ = img.shape
    return width, height

def add_image_to_slide(prs, img_path):
    # Add a slide with the image
    slide_layout = prs.slide_layouts[5]  # Using a blank slide layout
    slide = prs.slides.add_slide(slide_layout)

    # Get image dimensions
    img_width, img_height = get_image_size(img_path)

    # Calculate scaling factor to fit image within the slide without distortion
    max_width = Inches(8)
    max_height = Inches(6)
    
    width_ratio = max_width / img_width
    height_ratio = max_height / img_height
    scaling_factor = min(width_ratio, height_ratio)

    # Calculate the new width and height of the image
    new_width = img_width * scaling_factor
    new_height = img_height * scaling_factor

    # Calculate the position to center the image on the slide
    left = (Inches(10) - new_width) / 2
    top = (Inches(7.5) - new_height) / 2

    slide.shapes.add_picture(img_path, left, top, width=new_width, height=new_height)

def numerical_sort(value):
    return int(value.split('/')[-1].split('_')[-1].split('.')[0])
    




status = {}

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Allows all origins
    allow_credentials=True,
    allow_methods=["*"],  # Allows all methods
    allow_headers=["*"],  # Allows all headers
)


class YouTubeLink(BaseModel):
    link: str

class ROI(BaseModel):
    x: float
    y: float 
    width: float
    height: float
    res_width: int
    res_height: int
    video_id: str

def string_to_sha256(input_string):
    # Create a new sha256 hash object
    sha256_hash = hashlib.sha256()
    
    # Update the hash object with the bytes of the input string
    sha256_hash.update(input_string.encode('utf-8'))
    
    # Get the hexadecimal representation of the digest
    hex_digest = sha256_hash.hexdigest()
    
    return hex_digest



def download_video(link):
    youtubeObject = YouTube(link)
    youtubeObject = youtubeObject.streams.get_highest_resolution()
    input_string = youtubeObject.title
    video_id = string_to_sha256(input_string)
    video_file = f"./videos/{video_id}.mp4"
    try:
        if not os.path.exists(video_file):
            youtubeObject.download(filename=video_file)
    except:
        print("An error has occurred")
    print("Download is completed successfully")
    
    print(video_file)
    cap = cv2.VideoCapture(video_file)
    cap.set(1, 10)
    ret, frame = cap.read()
    print(ret)
    print(frame.shape)
    _, buffer = cv2.imencode('.jpg', frame)
    image = base64.b64encode(buffer).decode('utf-8')

    return image, video_id


def run_ai_algorithm(roi):

    x,y,w,h, video_id, res_width, res_height = roi.x, roi.y, roi.width, roi.height, roi.video_id, roi.res_width, roi.res_height



    os.makedirs(f'./frames/{video_id}', exist_ok=True)
    #counter = 1
    
    # Load the video
    cap = cv2.VideoCapture(f"./videos/{video_id}.mp4")
    fps = cap.get(cv2.CAP_PROP_FPS)
    frame_count = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))

    video_width = cap.get(cv2.CAP_PROP_FRAME_WIDTH)
    video_height = cap.get(cv2.CAP_PROP_FRAME_HEIGHT)

    x = int((x/res_width)*video_width)
    y = int((y/res_height)*video_height)
    w = int((w/res_width)*video_width)
    h = int((h/res_height)*video_height)
    # Define ROI (manually set these values based on your video)
    # x, y, w, h = 4,4,837,997  # Example values

    bbox = None
    fromCenter = False
    # Select multiple rectangles
    # selectROI("Image", im, rects, fromCenter)


    cropped_frames = []
    prev_features = None
    threshold = 5  # Define a threshold for feature difference (Euclidean distance)

    # Process every nth frame to speed up (e.g., every second frame)
    frame_interval = int(fps)  # Adjust this value to process fewer frames
    start_time = time.time()
    counter = 0
    for i in tqdm(range(0, frame_count, frame_interval)):
        cap.set(cv2.CAP_PROP_POS_FRAMES, i)
        ret, frame = cap.read()
        if not ret:
            break
        counter += 1
        cropped_frame = frame[y:y+h, x:x+w]

        # Extract features from the cropped frame
        current_features = extract_features(cropped_frame)
        
        if prev_features is not None:
            # Calculate the Euclidean distance between the current features and the previous features
            distance = np.linalg.norm(current_features - prev_features)
            
            # If the distance is above the threshold, consider it a new slide
            print(distance)
            if distance > threshold:
                # cropped_frames.append(cropped_frame)
                cv2.imwrite(f'./frames/{video_id}/crop_{i}.png', frame)
        else:
            # cropped_frames.append(cropped_frame)
            cv2.imwrite(f'./frames/{video_id}/crop_{i}.png', frame)
        
        prev_features = current_features
        print('FPS:', counter/(time.time()-start_time))

    cap.release()

def create_ppt(video_id):

    print(f'Creating ppt for {video_id}')
    # Create a presentation object
    prs = Presentation()

    # Get a sorted list of image filenames with numerical sorting
    image_files = sorted(glob(f'./frames/{video_id}/*.png'), key=numerical_sort)
    #print(image_files)
    # Iterate over sorted images and add to slides
    for img_filename in image_files:
        add_image_to_slide(prs, img_filename)

    # Save the presentation
    prs.save(f'./ppts/{video_id}.pptx')
    status[video_id] = 'processed'

processing_ready = False

@app.post("/api/youtube")
async def process_youtube_link(youtube_link: YouTubeLink, background_tasks: BackgroundTasks):

    image, video_id = download_video(youtube_link.link)
    
    return JSONResponse(content={"image": image, "video_id":video_id})

@app.post("/api/roi")
async def process_roi(roi: ROI, background_tasks: BackgroundTasks):

    if os.path.exists(f'./ppts/{roi.video_id}.pptx'):
        status[roi.video_id] = 'processed'
        return {"ready": True}
    else:
        status[roi.video_id] = 'processing'
        background_tasks.add_task(run_ai_algorithm, roi)
        background_tasks.add_task(create_ppt, roi.video_id)
        return {"ready": False}

@app.get("/api/checkReady/{video_id}")
async def check_ready(video_id:str):
    temp_status = status.get(video_id, "not found")
    return {"ready": True if temp_status == 'processed' else False, "downloadUrl": f"/get-ppt/{video_id}" }


@app.get("/get-ppt/{video_id}")
async def get_ppt(video_id:str):
    return FileResponse(f'./ppts/{video_id}.pptx')
